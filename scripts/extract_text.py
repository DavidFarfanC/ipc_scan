import os
import fitz  # PyMuPDF para PDF
import docx
import pandas as pd
import pptx
from bs4 import BeautifulSoup

def extract_text_from_pdf(file_path):
    """
    Extrae texto de un archivo PDF.
    """
    text = ""
    try:
        with fitz.open(file_path) as pdf:
            for page in pdf:
                text += page.get_text()
    except Exception as e:
        print(f"⚠️ Error al leer {file_path}: {e}")
    return text

def extract_text_from_docx(file_path):
    """
    Extrae texto de un archivo DOCX (Word).
    """
    text = ""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        print(f"⚠️ Error al leer {file_path}: {e}")
    return text

def extract_text_from_xlsx(file_path):
    """
    Extrae texto de un archivo XLSX (Excel).
    """
    text = ""
    try:
        df = pd.read_excel(file_path, engine="openpyxl")
        text = df.to_string()
    except Exception as e:
        print(f"⚠️ Error al leer {file_path}: {e}")
    return text

def extract_text_from_csv(file_path):
    """
    Extrae texto de un archivo CSV.
    """
    text = ""
    try:
        df = pd.read_csv(file_path, encoding="utf-8")
        text = df.to_string()
    except Exception as e:
        print(f"⚠️ Error al leer {file_path}: {e}")
    return text

def extract_text_from_pptx(file_path):
    """
    Extrae texto de un archivo PPTX (PowerPoint).
    """
    text = ""
    try:
        presentation = pptx.Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"⚠️ Error al leer {file_path}: {e}")
    return text

def extract_text_from_html(file_path):
    """
    Extrae texto de un archivo HTML.
    """
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
            text = soup.get_text()
    except Exception as e:
        print(f"⚠️ Error al leer {file_path}: {e}")
    return text

def extract_text_from_txt(file_path):
    """
    Extrae texto de un archivo TXT.
    """
    try:
        return open(file_path, "r", encoding="utf-8", errors="ignore").read()
    except Exception as e:
        print(f"⚠️ Error al leer {file_path}: {e}")
        return ""

def extract_text_from_script(file_path):
    """
    Extrae texto de archivos de código fuente (Python, Bash).
    """
    try:
        return open(file_path, "r", encoding="utf-8", errors="ignore").read()
    except Exception as e:
        print(f"⚠️ Error al leer {file_path}: {e}")
        return ""

def extract_text(file_path):
    """
    Identifica el tipo de archivo y extrae el texto en consecuencia.
    """
    ext = os.path.splitext(file_path)[-1].lower()
    
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file_path)
    elif ext == ".csv":
        return extract_text_from_csv(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".html":
        return extract_text_from_html(file_path)
    elif ext in [".txt", ".log"]:
        return extract_text_from_txt(file_path)
    elif ext in [".py", ".sh", ".bash"]:
        return extract_text_from_script(file_path)
    
    return None  # Si no es un tipo de archivo soportado
